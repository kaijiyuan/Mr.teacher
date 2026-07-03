"""Build .pptx files from PPTAgent's JSON outline.

Usage:
    from app.agents.ppt.builder import build_pptx

    file_path = build_pptx(outline_dict, output_dir="data/ppt")
"""

from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── colour palette ──────────────────────────────────────────────────────

PRIMARY = RGBColor(0x0F, 0x60, 0x4F)
SECONDARY = RGBColor(0x15, 0x6F, 0x5B)
DARK = RGBColor(0x1E, 0x23, 0x29)
BODY = RGBColor(0x4A, 0x55, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xF9)
ACCENT = RGBColor(0xE5, 0xF3, 0xEF)
MUTED = RGBColor(0x98, 0xA2, 0xA8)
GOLD = RGBColor(0xF0, 0xA0, 0x30)
DIVIDER = RGBColor(0xD8, 0xE0, 0xE2)


def build_pptx(
    outline: dict[str, Any],
    output_dir: str | Path = "data/ppt",
) -> str:
    """Generate a .pptx file from a PPTAgent outline dict.

    Args:
        outline: The ``ppt_outline`` dict from PPTAgent's state_update.
        output_dir: Directory to save the generated file.

    Returns:
        Absolute path to the generated .pptx file.
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = outline.get("slides", [])
    total = len(slides_data)

    # ── Title slide ───────────────────────────────────────────────────
    _add_title_slide(prs, outline.get("title", "复习 PPT"),
                     outline.get("summary", ""))

    # ── Content slides ────────────────────────────────────────────────
    current_topic = ""
    for i, slide_data in enumerate(slides_data, 1):
        # 检测主题切换，插入分隔页
        title = slide_data.get("title", "")
        topic_base = title.split("：")[0].split("——")[0].strip() if "：" in title or "——" in title else ""
        if topic_base and topic_base != current_topic and i > 1:
            current_topic = topic_base
        _add_content_slide(prs, slide_data, i, total, topic_hint=topic_base)

    # ── Closing slide ─────────────────────────────────────────────────
    _add_closing_slide(prs, outline.get("title", "复习 PPT"))

    # ── Save ──────────────────────────────────────────────────────────
    timestamp = datetime.now().strftime("%Y%m%d__%H%M%S")
    safe_title = outline.get("title", "untitled").lower().replace(" ", "_")[:32]
    filename = f"ppt__{safe_title}__{timestamp}.pptx"
    full_path = output_path / filename
    prs.save(str(full_path))
    return str(full_path.resolve())


# ═══════════════════════════════════════════════════════════════════════
#  Slide builders
# ═══════════════════════════════════════════════════════════════════════


def _add_title_slide(prs: Presentation, title: str, summary: str) -> None:
    """Professional title slide with decorative elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # ── Background ──
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY

    # ── Decorative top-right circle ──
    circle = slide.shapes.add_shape(
        9,  # oval
        w - Inches(2.5), Inches(-1.2),
        Inches(4.0), Inches(4.0),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x0D, 0x55, 0x46)
    circle.line.fill.background()

    # ── Decorative bottom bar ──
    bar = slide.shapes.add_shape(1, Inches(0), h - Inches(0.15), w, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    # ── Title ──
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), w - Inches(2.4), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # ── Subtitle line (gold accent line under title) ──
    line_shape = slide.shapes.add_shape(1, Inches(1.2), Inches(3.9), Inches(4), Inches(0.06))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = GOLD

    # ── Summary ──
    if summary:
        tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), w - Inches(3), Inches(2))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = summary
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xC8, 0xDD, 0xD8)
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = Pt(28)

    # ── Footer hint ──
    tb3 = slide.shapes.add_textbox(Inches(1.2), h - Inches(0.9), Inches(5), Inches(0.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "AI 智能学习助手 · 复习资料"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0x80, 0xB0, 0xA0)

    # Page number for title
    _add_page_number(slide, "1")


def _add_content_slide(
    prs: Presentation,
    data: dict[str, Any],
    index: int,
    total: int,
    topic_hint: str = "",
) -> None:
    """Content slide with left accent bar, styled bullets, and notes area hint."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # ── White background ──
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # ── Top decorative bar ──
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), w, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # ── Left accent strip ──
    accent = slide.shapes.add_shape(1, Inches(0.4), Inches(1.0), Inches(0.08), Inches(0.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SECONDARY
    accent.line.fill.background()

    # ── Title ──
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), w - Inches(2), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", f"第 {index} 页")
    p.font.size = Pt(30)
    p.font.color.rgb = DARK
    p.font.bold = True

    # ── Topic tag (if topic differs) ──
    if topic_hint:
        tag_tb = slide.shapes.add_textbox(w - Inches(2.5), Inches(0.35), Inches(2.2), Inches(0.35))
        tag_tf = tag_tb.text_frame
        tag_p = tag_tf.paragraphs[0]
        tag_p.text = topic_hint
        tag_p.font.size = Pt(10)
        tag_p.font.color.rgb = WHITE
        # Tag background is tricky in python-pptx, skip for simplicity

    # ── Content area (two-column if enough bullets) ──
    bullets = data.get("bullets", [])
    if bullets:
        max_lines = min(len(bullets), 14)
        col_count = 2 if max_lines >= 6 else 1
        items_per_col = (max_lines + col_count - 1) // col_count

        for col_idx in range(col_count):
            col_start = col_idx * items_per_col
            col_end = min(col_start + items_per_col, max_lines)
            if col_start >= col_end:
                break

            x_start = Inches(0.7) if col_idx == 0 else Inches(6.8)
            col_width = Inches(5.8)

            y_start = Inches(1.6)
            line_h = Inches(0.48)

            tb2 = slide.shapes.add_textbox(x_start, y_start, col_width, line_h * (col_end - col_start))
            tf2 = tb2.text_frame
            tf2.word_wrap = True

            for j, bullet in enumerate(bullets[col_start:col_end]):
                if j == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()

                # 检测是否为子要点（以 "- " 或 "  " 开头）
                is_sub = bullet.startswith("  ") or bullet.startswith("- ")
                display_text = bullet.lstrip(" -")

                if is_sub:
                    p2.text = f"    ›  {display_text}"
                    p2.font.size = Pt(16)
                    p2.font.color.rgb = MUTED
                else:
                    p2.text = f"  {display_text}"
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = DARK
                    p2.font.bold = False

                p2.space_after = Pt(6)
                p2.alignment = PP_ALIGN.LEFT

    # ── Right decorative circle (subtle) ──
    deco = slide.shapes.add_shape(9, w - Inches(1.2), h - Inches(1.2), Inches(1.0), Inches(1.0))
    deco.fill.solid()
    deco.fill.fore_color.rgb = ACCENT
    deco.line.fill.background()

    # ── Speaker notes ──
    notes = data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    # ── Bottom divider line ──
    divider = slide.shapes.add_shape(1, Inches(0.7), h - Inches(0.55), w - Inches(1.4), Inches(0.01))
    divider.fill.solid()
    divider.fill.fore_color.rgb = DIVIDER
    divider.line.fill.background()

    # ── Slide number + total info ──
    _add_page_number(slide, f"{index + 1} / {total + 2}")  # +2 for title + closing
    _add_slide_footer(slide, f"第 {index} 页")


def _add_closing_slide(prs: Presentation, title: str) -> None:
    """Professional closing / thank-you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # ── Background ──
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY

    # ── Decorative top circle ──
    circle = slide.shapes.add_shape(
        9, Inches(-1.0), Inches(-1.0), Inches(3.5), Inches(3.5),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x0D, 0x55, 0x46)
    circle.line.fill.background()

    # ── Decorative bottom bar ──
    bar = slide.shapes.add_shape(1, Inches(0), h - Inches(0.15), w, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    # ── "复习完成" text ──
    tb = slide.shapes.add_textbox(Inches(2), Inches(2.0), w - Inches(4), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "复习完成"
    p.font.size = Pt(50)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ── Topic name ──
    tb2 = slide.shapes.add_textbox(Inches(2), Inches(3.4), w - Inches(4), Inches(0.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xC8, 0xDD, 0xD8)
    p2.alignment = PP_ALIGN.CENTER

    # ── Encouraging message ──
    tb3 = slide.shapes.add_textbox(Inches(3), Inches(4.5), w - Inches(6), Inches(1))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "坚持复习，持续进步！"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(0xA0, 0xCC, 0xC0)
    p3.alignment = PP_ALIGN.CENTER

    # Page number
    _add_page_number(slide, f"{len(prs.slides)} / {len(prs.slides) + 1}")


# ═══════════════════════════════════════════════════════════════════════
#  Helpers
# ═══════════════════════════════════════════════════════════════════════


def _add_page_number(slide: Any, text: str) -> None:
    """Add page number at bottom-right."""
    tb = slide.shapes.add_textbox(
        Inches(11.0), Inches(7.0),
        Inches(2.0), Inches(0.4),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def _add_slide_footer(slide: Any, text: str) -> None:
    """Add small footer text at bottom-left."""
    tb = slide.shapes.add_textbox(
        Inches(0.7), Inches(7.0),
        Inches(4), Inches(0.4),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT
